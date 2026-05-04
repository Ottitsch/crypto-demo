"""
Generate a 16:9 PowerPoint presentation — Stage 1: Overall Approach & Initial Design.
Clean light palette matching academic style. Answers the 5 assignment questions only.

Run:
    python proposal/generate_presentation.py
"""
from __future__ import annotations
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Paths ──────────────────────────────────────────────────────────────────────
SCRIPT_DIR  = Path(__file__).resolve().parent
REPO_ROOT   = SCRIPT_DIR.parent
RESULTS_DIR = REPO_ROOT / "data" / "results"
ASSETS_DIR  = SCRIPT_DIR / "assets"
LOGO_PATH   = ASSETS_DIR / "tu_wien_informatics_logo.png"
OUT_PATH    = SCRIPT_DIR / "presentation.pptx"

# ── Palette: warm cream + charcoal + green accent ─────────────────────────────
# Background tones — warm cream instead of cold white
BG_WHITE    = RGBColor(0xFA, 0xF8, 0xF5)    # cream
BG_GRAY     = RGBColor(0xF3, 0xF0, 0xEB)    # warm light gray (hero sections)
CARD_BG     = RGBColor(0xF7, 0xF5, 0xF1)    # card fill — warm off-white
CARD_BORDER = RGBColor(0xE2, 0xDE, 0xD8)    # warm subtle border

# Text tones
INK         = RGBColor(0x1E, 0x1E, 0x1E)    # headings — near-black
BODY        = RGBColor(0x3A, 0x3A, 0x3A)    # body text — dark gray
MUTED       = RGBColor(0x6B, 0x6B, 0x6B)    # secondary / captions
LIGHT_TEXT  = RGBColor(0x9A, 0x9A, 0x9A)    # footer / very muted

# Accent — vibrant green
GREEN       = RGBColor(0x34, 0xC7, 0x7C)    # primary accent (sparingly)
GREEN_DARK  = RGBColor(0x22, 0x8B, 0x56)    # darker green for text on cream
GREEN_LIGHT = RGBColor(0xEA, 0xF7, 0xEF)    # very light green tint (highlight bg)

# Table header
TBL_HEAD_BG = RGBColor(0x2D, 0x2D, 0x2D)   # dark charcoal header row
TBL_ALT     = RGBColor(0xF5, 0xF3, 0xEF)   # warm alt row

# ── Fonts (matching Iknaio: Brawler for headings, General Sans for body) ──────
FONT_HEADING = "Brawler"       # serif — big titles
FONT_BODY    = "General Sans"  # geometric sans — everything else

# ── Slide dimensions (16:9) ───────────────────────────────────────────────────
W = Inches(13.333)
H = Inches(7.5)


def _hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


# ── XML helpers ───────────────────────────────────────────────────────────────

def _transition(slide) -> None:
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    node = slide.shapes._spTree.getparent()
    while node.tag != f"{{{p_ns}}}sld":
        node = node.getparent()
        if node is None:
            return
    for ch in list(node):
        if ch.tag == f"{{{p_ns}}}transition":
            node.remove(ch)
    node.append(etree.fromstring(
        f'<p:transition xmlns:p="{p_ns}" spd="med"><p:fade/></p:transition>'
    ))


def _bg(slide, color: RGBColor) -> None:
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    cSld = slide.shapes._spTree.getparent()
    old = cSld.find(f"{{{p_ns}}}bg")
    if old is not None:
        cSld.remove(old)
    cSld.insert(0, etree.fromstring(
        f'<p:bg xmlns:p="{p_ns}" xmlns:a="{a_ns}">'
        f'<p:bgPr><a:solidFill><a:srgbClr val="{_hex(color)}"/>'
        f'</a:solidFill><a:effectLst/></p:bgPr></p:bg>'
    ))


# ── Drawing primitives ────────────────────────────────────────────────────────

def _rect(sl, x, y, w, h, fill=None, line=None) -> None:
    s = sl.shapes.add_shape(1, x, y, w, h)
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if line:
        s.line.color.rgb = line
        s.line.width = Pt(0.5)
    else:
        s.line.fill.background()


def _txt(sl, text, x, y, w, h, sz=14, bold=False, italic=False,
         color=BODY, align=PP_ALIGN.LEFT, heading=False) -> None:
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT_HEADING if heading else FONT_BODY


def _bullets(sl, items, x, y, w, h, sz=12, color=BODY) -> None:
    tb = sl.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = item
        r.font.size = Pt(sz)
        r.font.color.rgb = color
        r.font.name = FONT_BODY
        p.space_before = Pt(4)
        p.space_after = Pt(2)


def _table(sl, headers, rows, x, y, w, h, col_widths=None) -> None:
    nc = len(headers)
    nr = len(rows) + 1
    if col_widths is None:
        col_widths = [w // nc] * nc
    tbl = sl.shapes.add_table(nr, nc, x, y, w, h).table
    for ci, cw in enumerate(col_widths):
        tbl.columns[ci].width = cw
    for ci, hd in enumerate(headers):
        c = tbl.cell(0, ci)
        c.text = hd
        p = c.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(10)
        p.font.color.rgb = BG_WHITE
        p.font.name = FONT_BODY
        c.fill.solid()
        c.fill.fore_color.rgb = TBL_HEAD_BG
    for ri, row in enumerate(rows):
        bg = TBL_ALT if ri % 2 == 0 else BG_WHITE
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.text = val
            p = c.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.bold = (ci == 0)
            p.font.color.rgb = INK if ci == 0 else BODY
            p.font.name = FONT_BODY
            c.fill.solid()
            c.fill.fore_color.rgb = bg


def _logo(sl) -> None:
    if LOGO_PATH.exists():
        sl.shapes.add_picture(
            str(LOGO_PATH), W - Inches(2.5), Inches(0.15),
            Inches(2.3), Inches(0.38),
        )


def _footer(sl, n, total=11) -> None:
    _rect(sl, 0, H - Inches(0.32), W, Inches(0.32), fill=BG_GRAY)
    _txt(sl, f"Crypto Asset Analytics  ·  TU Wien  ·  2026S",
         Inches(0.4), H - Inches(0.29), Inches(6), Inches(0.24),
         sz=8, color=LIGHT_TEXT)
    _txt(sl, f"{n} / {total}",
         W - Inches(1.0), H - Inches(0.29), Inches(0.7), Inches(0.24),
         sz=8, color=LIGHT_TEXT, align=PP_ALIGN.RIGHT)


def _header(sl, title, sub="") -> None:
    """Minimal header: thin green top-line + title on white."""
    _rect(sl, 0, 0, W, Inches(0.06), fill=GREEN)
    _logo(sl)
    _txt(sl, title, Inches(0.5), Inches(0.22), Inches(9.5), Inches(0.55),
         sz=24, bold=True, color=INK, heading=True)
    if sub:
        _txt(sl, sub, Inches(0.5), Inches(0.72), Inches(9.5), Inches(0.3),
             sz=11, color=MUTED)


def _card(sl, x, y, w, h, title, body,
          title_sz=13, body_sz=11, accent=GREEN) -> None:
    _rect(sl, x, y, w, h, fill=CARD_BG, line=CARD_BORDER)
    _rect(sl, x, y, Inches(0.06), h, fill=accent)
    _txt(sl, title, x + Inches(0.22), y + Inches(0.12),
         w - Inches(0.35), Inches(0.36),
         sz=title_sz, bold=True, color=INK, heading=True)
    _txt(sl, body, x + Inches(0.22), y + Inches(0.52),
         w - Inches(0.35), h - Inches(0.62),
         sz=body_sz, color=BODY)


def _note(sl, text) -> None:
    sl.notes_slide.notes_text_frame.text = text


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def _s01_title(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_GRAY)
    _transition(sl)

    # Thin green top edge
    _rect(sl, 0, 0, W, Inches(0.06), fill=GREEN)

    _txt(sl, "Proof of Funds",
         Inches(0.6), Inches(1.5), Inches(9.5), Inches(1.2),
         sz=48, bold=True, color=INK, heading=True)
    _txt(sl, "Bitcoin Risk Quantification for KYT / AML Compliance",
         Inches(0.6), Inches(2.8), Inches(9.0), Inches(0.5),
         sz=16, color=MUTED)

    # Green accent bar
    _rect(sl, Inches(0.6), Inches(3.4), Inches(3.5), Inches(0.05), fill=GREEN)

    _txt(sl, "Overall Approach & Initial Design",
         Inches(0.6), Inches(3.6), Inches(8), Inches(0.4),
         sz=13, italic=True, color=MUTED)

    _txt(sl, "Group 3  ·  Crypto Asset Analytics",
         Inches(0.6), Inches(4.5), Inches(9), Inches(0.35),
         sz=12, color=INK)
    _txt(sl, "Ottitsch Franz Stefan Jakob (12024717)\n"
         "Arora Aaradhaya (12534787)\n"
         "Slutu Serhii (12537831)",
         Inches(0.6), Inches(4.9), Inches(8), Inches(0.9),
         sz=11, color=MUTED)
    _txt(sl, "TU Wien Informatics  ·  Summer 2026",
         Inches(0.6), Inches(5.85), Inches(6), Inches(0.3),
         sz=10, italic=True, color=LIGHT_TEXT)

    if LOGO_PATH.exists():
        sl.shapes.add_picture(
            str(LOGO_PATH), W - Inches(3.2), H - Inches(1.1),
            Inches(2.8), Inches(0.47),
        )

    _note(sl, "Title slide. Introduce project name, team, course. ~15 seconds.")


def _s02_context(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Context & Motivation",
            "Why is this problem relevant?")
    _footer(sl, 2)

    _card(sl, Inches(0.5), Inches(1.15), Inches(12.3), Inches(1.55),
          "Regulatory requirement",
          "Cryptoasset service providers must quantify the risk of incoming Bitcoin transactions "
          "as part of Know-Your-Transaction (KYT) and Anti-Money Laundering (AML) procedures. "
          "A compliance analyst needs to answer: where did these coins come from, and how close "
          "are they to known illicit activity?")

    _card(sl, Inches(0.5), Inches(2.88), Inches(12.3), Inches(1.55),
          "The transparency gap",
          "Commercial solutions (Chainalysis, Elliptic, Crystal) provide risk scores, "
          "but their scoring logic is proprietary and unexplainable. An analyst cannot audit "
          "why a particular address was flagged — making regulatory justification difficult.",
          accent=RGBColor(0x2D, 0x2D, 0x2D))

    _card(sl, Inches(0.5), Inches(4.6), Inches(12.3), Inches(1.55),
          "Our goal",
          "Build an open-source, fully reproducible Bitcoin risk-scoring prototype "
          "where every score decomposes into auditable, interpretable components. "
          "We use the GraphSense TagPack infrastructure — an open standard for "
          "provenance-aware cryptoasset attribution — as our ground truth.",
          accent=GREEN)

    _note(sl, "Regulatory context → black-box problem → our open-source explainable approach. ~45 seconds.")


def _s03_questions(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Research Questions",
            "What do we want to answer?")
    _footer(sl, 3)

    qs = [
        ("Q1  ·  Risk Quantification",
         "How can we quantify the risk of an incoming Bitcoin transaction "
         "by its graph distance from known illicit entities — ransomware "
         "operators, darknet markets, exchange hacks?"),
        ("Q2  ·  Score Design",
         "How do we combine graph proximity, direct value exposure, and "
         "multi-hop taint propagation into a single explainable 0–100 "
         "aggregate score?"),
        ("Q3  ·  Entity Clustering",
         "How much does grouping addresses into real-world entities — via "
         "co-spend and change-address heuristics — change the risk landscape "
         "compared to raw address-level analysis?"),
        ("Q4  ·  Empirical Validation",
         "Can we validate our scores against external ground truth: OFAC sanctions "
         "and Bitcoin addresses from documented criminal cases?"),
    ]

    for i, (title, body) in enumerate(qs):
        col, row = i % 2, i // 2
        cx = Inches(0.5) + col * Inches(6.4)
        cy = Inches(1.15) + row * Inches(2.65)
        _card(sl, cx, cy, Inches(6.15), Inches(2.45), title, body)

    _note(sl, "Four research questions. ~15 s each. ~60 seconds total.")


def _s04_data(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Data Sources",
            "What data do we need and how do we get it?")
    _footer(sl, 4)

    _table(
        sl,
        headers=["Source", "Description", "Format", "Access"],
        rows=[
            ["GraphSense TagPacks",
             "Open address labels: ransomware, darknet markets, scams, hacks, exchanges (~46K addresses)",
             "YAML", "Free / GitHub"],
            ["Custom TagPacks",
             "Hand-curated from DOJ press releases for specific case studies",
             "YAML", "Manual / public"],
            ["mempool.space API",
             "Public blockchain explorer — provides full transaction data for neighbourhood crawl",
             "JSON REST", "Free / no key"],
            ["OFAC SDN List",
             "US Treasury sanctioned cryptocurrency addresses — serves as validation ground truth",
             "TXT", "Free / Gov"],
            ["Ransomwhere",
             "Crowdsourced ransomware payment addresses with family labels",
             "JSON", "Free / API"],
        ],
        x=Inches(0.5), y=Inches(1.15), w=Inches(12.3), h=Inches(3.5),
        col_widths=[
            int(Inches(2.5)), int(Inches(5.55)),
            int(Inches(1.25)), int(Inches(3.0)),
        ],
    )

    # Key principle box
    _rect(sl, Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.7),
          fill=GREEN_LIGHT, line=GREEN)
    _txt(sl,
         "Design principle:  all data sources are free, require no API keys, "
         "and API responses will be cached in SQLite — making the pipeline "
         "fully reproducible offline by anyone.",
         Inches(0.7), Inches(4.92), Inches(11.9), Inches(0.55),
         sz=11, color=GREEN_DARK)

    _note(sl, "Walk through five data sources. Emphasise reproducibility: no paid services, no API keys, fully open. ~60 seconds.")


def _s05_approach(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Analytics Approach",
            "How do we want to approach the research questions?")
    _footer(sl, 5)

    # Five-step pipeline as horizontal cards
    steps = [
        ("1", "Label Risk Sources",
         "Parse GraphSense TagPack YAML files to build a labelled "
         "set of seed addresses with category tags and severity weights."),
        ("2", "Build Transaction Graph",
         "BFS-crawl the transaction neighbourhood of seed addresses via "
         "mempool.space. Construct a directed, value-weighted graph where "
         "nodes are addresses and edges encode proportional value flow."),
        ("3", "Cluster Entities",
         "Apply Union-Find with two standard Bitcoin forensics heuristics "
         "(common-input-ownership + change-address detection) to group "
         "addresses into real-world entities."),
        ("4", "Compute Risk Scores",
         "Calculate four complementary metrics per node: BFS distance, "
         "direct exposure, haircut taint propagation, and a weighted 0–100 "
         "aggregate. Each component is independently interpretable."),
        ("5", "Validate & Analyze",
         "Evaluate scores against OFAC sanctions and documented criminal "
         "case addresses. Analyze score distributions and taint decay "
         "patterns to answer our research questions."),
    ]

    for i, (num, title, body) in enumerate(steps):
        cy = Inches(1.1) + i * Inches(1.18)
        # Number badge
        _rect(sl, Inches(0.5), cy, Inches(0.42), Inches(0.42), fill=GREEN)
        _txt(sl, num, Inches(0.5), cy + Inches(0.02),
             Inches(0.42), Inches(0.38),
             sz=14, bold=True, color=BG_WHITE, align=PP_ALIGN.CENTER)
        # Title
        _txt(sl, title, Inches(1.05), cy + Inches(0.04),
             Inches(3.0), Inches(0.36),
             sz=12, bold=True, color=INK, heading=True)
        # Body
        _txt(sl, body, Inches(4.2), cy + Inches(0.0),
             Inches(8.6), Inches(1.08),
             sz=11, color=BODY)
        # Separator
        if i < 4:
            _rect(sl, Inches(0.5), cy + Inches(1.08),
                  Inches(12.3), Inches(0.01), fill=CARD_BORDER)

    _note(sl, "Five-stage pipeline. Walk through each stage in one sentence. This slide answers 'analytics approach'. ~75 seconds.")


def _s06_metrics(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Risk Scoring Design",
            "Four complementary metrics — each independently interpretable")
    _footer(sl, 6)

    metrics = [
        ("BFS Distance",
         "Multi-source BFS from all tainted seeds. Each node receives the "
         "shortest hop-count to its nearest illicit predecessor.",
         "Proximity signal — fewer hops means more direct association."),
        ("Direct Exposure",
         "Fraction of 1-hop incoming value originating from tainted "
         "predecessors. Pure value-flow signal, independent of graph topology.",
         "How much incoming money comes directly from illicit sources."),
        ("Haircut Taint",
         "Multi-hop proportional propagation via sparse-matrix Jacobi iteration. "
         "Taint dilutes proportionally at each hop based on value fractions.",
         "Captures indirect contamination along extended laundering paths."),
        ("Aggregate Score",
         "Weighted blend:  0.30 × Distance  +  0.30 × Exposure  +  0.40 × Haircut, "
         "rescaled to 0–100.",
         "Single actionable score. Designed to be decomposable for compliance audit."),
    ]

    for i, (title, method, interp) in enumerate(metrics):
        col, row = i % 2, i // 2
        cx = Inches(0.5) + col * Inches(6.4)
        cy = Inches(1.15) + row * Inches(2.65)
        bw, bh = Inches(6.15), Inches(2.45)
        _rect(sl, cx, cy, bw, bh, fill=CARD_BG, line=CARD_BORDER)
        _rect(sl, cx, cy, Inches(0.06), bh,
              fill=GREEN if i < 3 else RGBColor(0x2D, 0x2D, 0x2D))
        _txt(sl, title, cx + Inches(0.2), cy + Inches(0.12),
             bw - Inches(0.3), Inches(0.34),
             sz=13, bold=True, color=INK, heading=True)
        _txt(sl, method, cx + Inches(0.2), cy + Inches(0.52),
             bw - Inches(0.3), Inches(1.1),
             sz=11, color=BODY)
        _rect(sl, cx + Inches(0.2), cy + Inches(1.68),
              bw - Inches(0.4), Inches(0.01), fill=CARD_BORDER)
        _txt(sl, interp, cx + Inches(0.2), cy + Inches(1.78),
             bw - Inches(0.3), Inches(0.55),
             sz=10, italic=True, color=MUTED)

    _note(sl, "Four metrics. Stress: the aggregate formula is public and each component has independent meaning. ~60 seconds.")


def _s07_clustering(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Entity Clustering",
            "Grouping addresses into real-world actors before scoring")
    _footer(sl, 7)

    # Heuristics
    _txt(sl, "Two standard Bitcoin forensics heuristics (Meiklejohn et al. 2013)",
         Inches(0.5), Inches(1.15), Inches(10), Inches(0.3),
         sz=11, bold=True, color=MUTED)

    _card(sl, Inches(0.5), Inches(1.55), Inches(5.9), Inches(2.1),
          "H1  ·  Common-Input-Ownership",
          "All input addresses in the same Bitcoin transaction belong to the same entity. "
          "Rationale: spending multiple UTXOs requires holding all corresponding private keys "
          "— a strong indicator of common ownership. This is the most widely used clustering "
          "heuristic in blockchain forensics.")

    _card(sl, Inches(6.9), Inches(1.55), Inches(5.9), Inches(2.1),
          "H2  ·  Change-Address Detection",
          "A freshly generated, single-use output receiving unspent change is grouped with "
          "the spending inputs. Standard wallet software returns leftover value to a new "
          "self-controlled address. We detect these and merge them with the input cluster.")

    # Why this matters
    _rect(sl, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.06), fill=GREEN)
    _txt(sl, "Why this matters for risk scoring",
         Inches(0.5), Inches(4.1), Inches(10), Inches(0.3),
         sz=12, bold=True, color=INK, heading=True)
    _txt(sl,
         "Without clustering, risk scores apply to individual addresses. But in Bitcoin, a single "
         "user may control hundreds of addresses. Clustering consolidates the risk signal — "
         "revealing the true exposure of a real-world entity rather than fragmenting it across "
         "unrelated-looking addresses. We plan to measure how much clustering changes the "
         "overall risk distribution (our research question Q3).",
         Inches(0.5), Inches(4.45), Inches(12.3), Inches(1.2),
         sz=11, color=BODY)

    # Implementation note
    _txt(sl, "Implementation: Union-Find (disjoint-set) data structure — "
         "near-linear time complexity; scalable to large graphs.",
         Inches(0.5), Inches(5.75), Inches(12.3), Inches(0.3),
         sz=10, italic=True, color=MUTED)

    _note(sl, "Two heuristics from Meiklejohn 2013 — the foundational paper for Bitcoin forensics. Explain why clustering matters for risk. ~50 seconds.")


def _s08_implementation(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Implementation Plan",
            "How are we going to build this?")
    _footer(sl, 8)

    # Three-column architecture
    cols = [
        ("Data Layer",
         [
             "GraphSense TagPack parser with category weights",
             "Bounded BFS crawler with rate limiting",
             "SQLite response cache for reproducibility",
             "Entity label aggregation (46K+ addresses)",
         ]),
        ("Algorithm Layer",
         [
             "Union-Find clustering (H1 + H2 heuristics)",
             "Multi-source BFS distance computation",
             "Proportional exposure calculation",
             "Sparse-matrix Jacobi taint propagation",
             "Weighted aggregate scoring (0–100)",
         ]),
        ("Output & Validation Layer",
         [
             "Per-address and per-entity risk table",
             "Score distribution analysis",
             "ROC / AUC validation vs. OFAC ground truth",
             "Interactive investigation dashboard",
             "Jupyter notebooks for narrative presentation",
         ]),
    ]

    for ci, (col_title, items) in enumerate(cols):
        cx = Inches(0.5) + ci * Inches(4.25)
        cw = Inches(4.0)
        _rect(sl, cx, Inches(1.15), cw, Inches(4.55),
              fill=CARD_BG, line=CARD_BORDER)
        _rect(sl, cx, Inches(1.15), cw, Inches(0.05), fill=GREEN)
        _txt(sl, col_title,
             cx + Inches(0.2), Inches(1.28), cw - Inches(0.3), Inches(0.32),
             sz=12, bold=True, color=INK, heading=True)
        for ii, item in enumerate(items):
            _txt(sl, "›  " + item,
                 cx + Inches(0.2), Inches(1.72) + ii * Inches(0.62),
                 cw - Inches(0.3), Inches(0.55),
                 sz=10, color=BODY)

    # Tech stack — one line at bottom
    _txt(sl, "Stack:  Python 3.10+  ·  NumPy / SciPy  ·  NetworkX  ·  Streamlit  ·  pytest",
         Inches(0.5), Inches(5.85), Inches(12.3), Inches(0.28),
         sz=10, italic=True, color=MUTED)

    _note(sl, "Three-layer architecture. Don't read the lists — just name each layer and its purpose. Mention Streamlit dashboard briefly. ~60 seconds.")


def _s09_cases(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Planned Case Studies",
            "Real criminal cases with publicly known addresses from DOJ filings")
    _footer(sl, 9)

    _table(
        sl,
        headers=["Case", "Year", "Incident", "Why we chose it"],
        rows=[
            ["WannaCry",
             "2017",
             "Global ransomware — 200K+ victims in 150 countries",
             "Isolated wallets, minimal mixing — baseline for scoring"],
            ["Twitter Hack",
             "2020",
             "130 VIP accounts hijacked for a BTC doubling scam",
             "Aggressive fund consolidation — stresses clustering heuristics"],
            ["Colonial Pipeline",
             "2021",
             "DarkSide ransomware shut US fuel supply for 6 days",
             "Affiliate model distributes funds — tests label coverage"],
            ["Bitfinex Hack",
             "2016",
             "119K BTC stolen; largest DOJ crypto seizure in 2022",
             "Peel-chain laundering — tests multi-hop taint propagation"],
        ],
        x=Inches(0.5), y=Inches(1.15), w=Inches(12.3), h=Inches(2.72),
        col_widths=[
            int(Inches(1.8)), int(Inches(0.6)),
            int(Inches(4.3)), int(Inches(5.6)),
        ],
    )

    _txt(sl, "Validation strategy",
         Inches(0.5), Inches(4.1), Inches(5), Inches(0.3),
         sz=12, bold=True, color=INK, heading=True)
    _bullets(sl, [
        "ROC / AUC — aggregate score as binary classifier "
        "(positives = OFAC-listed or TagPack-labelled addresses)",
        "Precision / Recall at practical thresholds — "
        "does the score rank known-bad addresses above the rest?",
        "Address-level vs. entity-level comparison — "
        "quantify how much clustering changes the risk distribution",
    ], Inches(0.5), Inches(4.45), Inches(12.3), Inches(1.6),
        sz=11, color=BODY)

    _note(sl, "One sentence per case — explain why each case stresses a different pipeline component. ~65 seconds.")


def _s10_outcome(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_WHITE)
    _transition(sl)
    _header(sl, "Expected Outcome",
            "What will the project produce?")
    _footer(sl, 10)

    outcomes = [
        ("Reproducible risk pipeline",
         "A Python package that takes any Bitcoin address, crawls its neighbourhood, "
         "clusters entities, and outputs a scored risk table. Fully reproducible "
         "from raw YAML tag files — no paid services, no API keys."),
        ("Explainable risk scores",
         "Every score decomposes into distance, direct exposure, and haircut taint. "
         "A compliance analyst can inspect exactly which component drives the risk "
         "and trace it back to specific graph paths."),
        ("Empirical answers to our research questions",
         "Quantitative evidence on how risk decays with graph distance, how "
         "entity clustering reshapes risk distributions, and whether open-source "
         "scoring can match the separation achieved by commercial tools."),
        ("Validation on real cases",
         "Documented comparison against OFAC sanctions and DOJ-traced addresses "
         "from four high-profile criminal investigations — not synthetic data, "
         "but actual law-enforcement outcomes."),
    ]

    for i, (title, body) in enumerate(outcomes):
        col, row = i % 2, i // 2
        cx = Inches(0.5) + col * Inches(6.4)
        cy = Inches(1.15) + row * Inches(2.55)
        _card(sl, cx, cy, Inches(6.15), Inches(2.35), title, body)

    _note(sl, "Four expected deliverables / outcomes. Stress validation against real criminal cases. ~55 seconds.")


def _s11_closing(prs):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(sl, BG_GRAY)
    _transition(sl)

    # Green top bar
    _rect(sl, 0, 0, W, Inches(0.06), fill=GREEN)

    _txt(sl, "Summary & Next Steps",
         Inches(0.6), Inches(0.35), Inches(9), Inches(0.52),
         sz=22, bold=True, color=INK, heading=True)

    if LOGO_PATH.exists():
        sl.shapes.add_picture(
            str(LOGO_PATH), W - Inches(3.0), Inches(0.2),
            Inches(2.7), Inches(0.45),
        )

    _rect(sl, Inches(0.6), Inches(0.9), Inches(4.5), Inches(0.02), fill=GREEN)

    # Summary bullets
    _txt(sl, "What we plan to build",
         Inches(0.6), Inches(1.1), Inches(6), Inches(0.3),
         sz=12, bold=True, color=INK, heading=True)
    _bullets(sl, [
        "End-to-end open-source Bitcoin KYT / AML pipeline",
        "Built on the GraphSense TagPack standard for attribution",
        "Entity clustering via Union-Find (H1 + H2 heuristics)",
        "Explainable 0–100 risk score with three interpretable components",
        "Validated against OFAC and four real DOJ criminal cases",
    ], Inches(0.6), Inches(1.48), Inches(6.5), Inches(2.2),
        sz=11, color=BODY)

    # Timeline
    _txt(sl, "Work plan",
         Inches(0.6), Inches(3.7), Inches(6), Inches(0.3),
         sz=12, bold=True, color=INK, heading=True)

    phases = [
        "Data acquisition & label loading — TagPacks, seeds, BFS crawl setup",
        "Core metrics implementation — distance, exposure, haircut, aggregate",
        "Entity clustering & case study pipeline",
        "Validation, dashboard, notebooks, and final report",
    ]
    _bullets(sl, phases,
             Inches(0.6), Inches(4.05), Inches(6.5), Inches(1.7),
             sz=11, color=BODY)

    # Right side: team
    _txt(sl, "Team",
         Inches(8.0), Inches(1.1), Inches(5), Inches(0.3),
         sz=12, bold=True, color=INK, heading=True)
    members = [
        ("Ottitsch Franz Stefan Jakob", "12024717"),
        ("Arora Aaradhaya", "12534787"),
        ("Slutu Serhii", "12537831"),
    ]
    for i, (name, uid) in enumerate(members):
        cy = Inches(1.5) + i * Inches(0.65)
        _rect(sl, Inches(8.0), cy, Inches(4.8), Inches(0.55),
              fill=CARD_BG, line=CARD_BORDER)
        _txt(sl, name,
             Inches(8.15), cy + Inches(0.06), Inches(3.6), Inches(0.24),
             sz=11, bold=True, color=INK)
        _txt(sl, uid,
             Inches(8.15), cy + Inches(0.3), Inches(3.6), Inches(0.2),
             sz=9, color=MUTED)

    # Thank you
    _rect(sl, Inches(0.6), Inches(6.0), Inches(12.2), Inches(0.02), fill=CARD_BORDER)
    _txt(sl, "Thank you  ·  Questions?",
         Inches(0.6), Inches(6.15), Inches(8), Inches(0.48),
         sz=20, bold=True, color=INK, heading=True)

    _footer(sl, 11)
    _note(sl, "Brief summary. Mention timeline. Thank audience, invite questions. ~30 seconds.")


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    _s01_title(prs)
    _s02_context(prs)
    _s03_questions(prs)
    _s04_data(prs)
    _s05_approach(prs)
    _s06_metrics(prs)
    _s07_clustering(prs)
    _s08_implementation(prs)
    _s09_cases(prs)
    _s10_outcome(prs)
    _s11_closing(prs)

    prs.save(str(OUT_PATH))
    print(f"Saved -> {OUT_PATH}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
